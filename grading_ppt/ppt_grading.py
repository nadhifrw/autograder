import os

# from lxml.etree import parse_xml
from lxml.etree import fromstring as parse_xml
from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT

from checking_folder.checking_files import check_files_in_folder
from grading_ppt.finding_pptx import find_pptx


def ppt_slides(path):
    pptx_list = find_pptx(path)
    slide_scores = {}  # Initialize slide score
    for pptx in pptx_list:
        dirname = os.path.basename(os.path.dirname(pptx))

        # print(f"Processing {pptx}")
        prs = Presentation(pptx)
        score = 0
        slide_lengths = len(prs.slides)

        # Access slide master (which contains the theme info)

        if slide_lengths >= 8:
            score = +10
            print(
                f"Presentation {pptx} has {slide_lengths} slides, which is sufficient."
            )
        # print(f"Slide lengths in {pptx}: {slide_lengths}")
        slide_scores[dirname] = score
        print(f"Scoring for {dirname} completed with score: {score}")

    return slide_scores  # Return the slide scores dictionary


# This function checks the theme that is being used in the pptx files
def checking_theme(path):
    pptx_list = find_pptx(path)
    theme_scores = {}  # Initialize slide score
    for pptx in pptx_list:
        dirname = os.path.basename(os.path.dirname(pptx))

        # print(f"Processing {pptx}")
        prs = Presentation(pptx)
        score = 0

        # Get the Slide Master and its part
        slide_master = prs.slide_master
        slide_master_part = slide_master.part

        # Get the Theme part
        theme_part = slide_master_part.part_related_by(RT.THEME)
        theme = parse_xml(theme_part.blob)  # theme here is an <a:theme> element
        # print(f"Theme part for {dirname}: {theme_part}")
        # print(f"Theme XML for {dirname}: {theme}")
        theme_name = theme.get("name")
        # print(f"Theme name: {theme_name}")

        # this block of code is not being used as it is quite unecessary to check each slide's theme
        # but can be use if needed
        # themes_found = {}
        #
        # for i, slide in enumerate(prs.slides):
        #     slide_layout = slide.slide_layout
        #     slide_master = slide_layout.slide_master
        #     theme_part = slide_master.part.part_related_by(RT.THEME)
        #     theme = fromstring(theme_part.blob)
        #
        #     theme_name = theme.get("name")
        #
        #     if theme_name not in themes_found:
        #         themes_found[theme_name] = []
        #     themes_found[theme_name].append(i + 1)
        #
        # print("Themes used:")
        # for theme_name, slide_numbers in themes_found.items():
        #     print(f"'{theme_name}': Slides {slide_numbers}")

        # Making sure the theme that is being used is not the default theme (white blank theme)
        if theme_name not in ["Office Theme", "Office 2013 - 2022 Theme"]:
            score = +10
            print(f"Presentation {pptx} uses a custom theme: {theme_name}.")

        theme_scores[dirname] = score
        print(f"Scoring for {dirname}'s theme completed with score: {score}")

    return theme_scores  # Return the theme scores dictionary


def checking_animation_tranisition(path):
    # This function is not used in the current implementation
    # It can be used to check animations and transitions in a specific way if needed
    pass


def checking_hyperlinks(path):
    # This function is not used in the current implementation
    # It can be used to check hyperlinks in a specific way if needed
    pass


def checking_action_button(path):
    # This function is not used in the current implementation
    # It can be used to check action buttons in a specific way if needed
    pass


def find_ppsx(path):
    # Get all files from the directory
    files = check_files_in_folder(path)

    ppsx_scores = {}

    if not files:
        return ppsx_scores

    ppsx_files = [
        f
        for f in files
        if f.endswith(".ppsx")
        and "ppt" in f.lower()
        and "show" in f.lower()
        and "-" in f
    ]

    for ppsx in ppsx_files:
        dirname = os.path.basename(
            os.path.dirname(ppsx)
        )  # Get folder name from each file
        filename = os.path.basename(ppsx)
        print(f"Found ppsx file: {filename}")

        score = 0
        if "show" in filename.lower():
            score += 10
            print(f"Found ppsx file with 'show' in the name: {filename}")
        else:
            print(f"No 'show' found in the ppsx file name: {filename}")

        ppsx_scores[dirname] = score
        print(f"Scoring for {dirname} completed with score: {score}")

    return ppsx_scores
